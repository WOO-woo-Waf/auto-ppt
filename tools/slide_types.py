from pptx.enum.text import PP_ALIGN
from pptx.util import Inches
from pptx.dml.color import RGBColor
from tools.font_utils import set_font
from tools.insert_placeholder import insert_placeholder_box
from tools.insert_image_file import find_matching_image, insert_image
from tools.tag_utils import (
    parse_tags,
    get_illustration_keyword,
    get_color_from_tags,
    is_emphasis,
)


def add_title_slide(prs, slide_data, style):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = slide_data["title"]
    subtitle = slide.placeholders[1]
    subtitle.text = slide_data.get("subtitle", "")
    set_font(slide.shapes.title.text_frame.paragraphs[0], style, path="title.title_text")
    set_font(subtitle.text_frame.paragraphs[0], style, path="title.subtitle_text")
    return slide


def add_section_slide(prs, slide_data, style):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    box_width = Inches(10)
    box_height = Inches(1.5)
    left = int((slide_width - box_width) / 2)
    top = int((slide_height - box_height) / 2)

    textbox = slide.shapes.add_textbox(left, top, box_width, box_height)
    tf = textbox.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = slide_data["title"]
    set_font(p, style, path="section.title_text")
    p.alignment = PP_ALIGN.CENTER
    return slide


def add_content_slide(prs, slide_data, style, image_dir="images"):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = slide_data["title"]
    set_font(slide.shapes.title.text_frame.paragraphs[0], style, path="content.title_text")

    tf = slide.placeholders[1].text_frame
    tf.clear()

    for bullet in slide_data["bullets"]:
        text, tags = parse_tags(bullet)
        p = tf.add_paragraph()
        p.text = text
        set_font(p, style, path="content.bullet_text")

        run = p.runs[0]
        color_hex = get_color_from_tags(tags)
        if color_hex:
            run.font.color.rgb = RGBColor(
                int(color_hex[1:3], 16),
                int(color_hex[3:5], 16),
                int(color_hex[5:7], 16),
            )
        if is_emphasis(tags):
            run.font.bold = True

        keyword = get_illustration_keyword(tags, text)
        if keyword:
            image_path = find_matching_image(keyword, image_dir)
            if image_path:
                insert_image(slide, image_path)
            else:
                insert_placeholder_box(slide, keyword)

    return slide


def add_closing_slide(prs, slide_data, style):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = slide_data["title"]
    tf = slide.placeholders[1].text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = slide_data.get("subtitle", "")
    set_font(slide.shapes.title.text_frame.paragraphs[0], style, path="closing.title_text", alignment="center")
    set_font(p, style, path="closing.subtitle_text")
    return slide
