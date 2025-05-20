from pptx.dml.color import RGBColor
from pptx.util import Inches
import os


def set_background(slide, background_config, prs):
    fill = slide.background.fill
    bg_type = background_config.get("type", "solid")

    if bg_type == "solid":
        color = background_config.get("color", "#FFFFFF")
        fill.solid()
        fill.fore_color.rgb = RGBColor(int(color[1:3], 16), int(color[3:5], 16), int(color[5:7], 16))

    elif bg_type == "image":
        image_path = background_config.get("image")
        if image_path and os.path.exists(image_path):
            width = prs.slide_width
            height = prs.slide_height
            pic = slide.shapes.add_picture(image_path, Inches(0), Inches(0), width=width, height=height)
            # 将背景图元素移到最底层
            sp_tree = slide.shapes._spTree
            sp_tree.remove(pic._element)
            sp_tree.insert(2, pic._element)  # 通常索引 0 是 placeholder，1 是背景，应插入靠前但不覆盖内容
    elif bg_type == "mixed":
        # 先设置底色再添加图像覆盖（如渐变图）
        color = background_config.get("color", "#FFFFFF")
        image_path = background_config.get("image")
        fill.solid()
        fill.fore_color.rgb = RGBColor(int(color[1:3], 16), int(color[3:5], 16), int(color[5:7], 16))
        if image_path and os.path.exists(image_path):
            width = prs.slide_width
            height = prs.slide_height
            pic = slide.shapes.add_picture(image_path, Inches(0), Inches(0), width=width, height=height)
            # 将背景图元素移到最底层
            sp_tree = slide.shapes._spTree
            sp_tree.remove(pic._element)
            sp_tree.insert(2, pic._element)  # 通常索引 0 是 placeholder，1 是背景，应插入靠前但不覆盖内容
