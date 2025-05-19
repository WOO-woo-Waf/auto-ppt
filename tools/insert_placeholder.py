from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import re


def extract_illustration_keyword(text):
    match = re.search(r"【插图】(.*)", text)
    return match.group(1).strip() if match else None


def insert_placeholder_box(slide, keyword):
    if not keyword:
        return
    left = Inches(5.5)
    top = Inches(2)
    width = Inches(3.5)
    height = Inches(2.5)
    shape = slide.shapes.add_shape(
        autoshape_type_id=1,  # Rectangle
        left=left,
        top=top,
        width=width,
        height=height
    )
    shape.text = f"[待添加插图：{keyword}]"
    text_frame = shape.text_frame
    p = text_frame.paragraphs[0]
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = RGBColor(200, 200, 200)
